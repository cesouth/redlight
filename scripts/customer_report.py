#!/usr/bin/env python3
"""Produce a customer-ready congestion study from raw GPS and a road network.

End-to-end pipeline for GPS data that carries **no speed column** -- position,
time, a mover id and a horizontal-accuracy value are enough:

    load network (GeoJSON)
      -> load points (tz-aware, local clock)
      -> HMM / Viterbi map matching (Newson & Krumm)
      -> derive_speeds from on-road displacement, using per-point accuracy
      -> dwell-aware + robust cleaning
      -> assign overall / peak / off-peak speeds onto the network
      -> analytics (temporal, day-type, congestion vs posted limit, structure)
      -> the deliverable, in one of three formats

Formats (``--format``, inferred from ``--out``'s extension when omitted):

``pdf`` (default)
    A fixed, print-ready multi-page document. No dependencies beyond
    matplotlib. Not editable.
``pptx``
    An editable 16:9 PowerPoint deck, one slide per section, with the
    supporting numbers as **native tables** rather than pictures of tables --
    the format to choose when the customer will restyle it or fold it into
    their own template. Requires ``pip install python-pptx``.
``png``
    A numbered PNG per section plus ``captions.md``, for dropping into your own
    template or deck. ``--out`` names the directory; ``--dpi`` sets resolution.

All three run the same analysis and share the same figures and section
headings, so they cannot disagree with each other. The look commits
deliberately to light, print-oriented styling rather than following the
reader's OS theme -- the raster maps cannot re-theme, and a document that
prints predictably is worth more here than one that adapts to a screen.

Usage
-----
    python scripts/customer_report.py \\
        --network roads.geojson \\
        --points  probes.csv \\
        --accuracy-col accuracy \\
        --tz America/New_York \\
        --title "Route 1 Corridor Study" \\
        --customer "Example County DOT" \\
        --out report.pdf           # or report.pptx, or --format png --out figs/

Requires the ``mapping`` extra for figures: ``pip install roadtraffic[mapping]``
(matplotlib). Every stage degrades gracefully -- a study area with no posted
speed limits, or GPS that never sampled a weekend, drops the affected section
and says so in "Data notes" rather than failing or inventing a number.
"""
from __future__ import annotations

import argparse
import sys
import textwrap
import warnings
from datetime import datetime, timezone

import numpy as np
import pandas as pd

import roadtraffic as rt

# --------------------------------------------------------------------------- #
# Palette. Light-mode values from the project's validated data-viz palette.
# Colour choices that matter, and why they are not the conventional ones:
#
#   * Absolute speed is a *magnitude*, so it gets a single-hue sequential ramp
#     (blue, light->dark). The familiar red-yellow-green traffic ramp is a
#     rainbow -- it puts a hue at the midpoint, and red-vs-green is precisely
#     the pair that deuteranopic and protanopic readers cannot separate. On a
#     printed page handed to a room, that is a real failure, not a theoretical
#     one.
#   * The congestion ratio *does* have a natural midpoint (1.0 = running at the
#     posted limit), so it gets a true diverging scale: red below, neutral grey
#     at, blue above. Two hues that read as opposite, and nothing at the middle.
#   * Map line-work starts at ramp step 250 rather than 100: the palest
#     sequential steps are designed to recede into the surface, which is right
#     for a filled choropleth and wrong for a 2px line that then vanishes.
# --------------------------------------------------------------------------- #
INK = "#0b0b0b"
INK_2 = "#52514e"
MUTED = "#898781"
GRID = "#e1e0d9"
BASELINE = "#c3c2b7"
SURFACE = "#fcfcfb"
PLANE = "#f9f9f7"
SERIES_1 = "#2a78d6"     # blue
SERIES_2 = "#eb6834"     # orange
GOOD = "#0ca30c"
WARNING = "#fab219"
CRITICAL = "#d03b3b"

# Sequential blue, steps 250..700 -- see note above on why not 100.
SEQ_BLUE = ["#86b6ef", "#5598e7", "#3987e5", "#2a78d6", "#256abf",
            "#1c5cab", "#184f95", "#104281", "#0d366b"]
# Diverging poles + neutral midpoint. The midpoint is the *baseline* grey, not
# the palette's pale diverging grey: on a near-white surface a pale midpoint
# makes every road near the middle of the scale disappear, and on this map the
# middle is where most roads live. A mid grey still reads as "neither pole"
# (it carries no hue) while staying legible as a 2px line.
DIV_LOW, DIV_MID, DIV_HIGH = "#d03b3b", "#c3c2b7", "#2a78d6"

# Categorical: three unordered classes. Distinct hues, not a ramp.
# Slots 1-3 of the validated categorical theme; these three clear the
# all-pairs CVD and normal-vision floors on this deck's surface.
CAT_1, CAT_2, CAT_3 = "#2a78d6", "#eb6834", "#1baf7a"
MODE_COLORS = {"vehicle": CAT_1, "pedestrian": CAT_2, "unknown": CAT_3}
MODE_ORDER = ("vehicle", "pedestrian", "unknown")

FONT_STACK = ["system-ui", "-apple-system", "Segoe UI", "Helvetica", "Arial",
              "DejaVu Sans", "sans-serif"]


def _mpl():
    """Import and configure matplotlib, with an actionable error if absent."""
    try:
        import matplotlib
    except ImportError as exc:  # pragma: no cover - depends on user's env
        raise SystemExit(
            "This report needs matplotlib for its figures.\n"
            "    pip install 'roadtraffic[mapping]'"
        ) from exc
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    plt.rcParams.update({
        "font.family": "sans-serif",
        "font.sans-serif": FONT_STACK,
        "figure.facecolor": SURFACE,
        "axes.facecolor": SURFACE,
        "axes.edgecolor": BASELINE,
        "axes.labelcolor": INK_2,
        "axes.titlecolor": INK,
        "text.color": INK,
        "xtick.color": MUTED,
        "ytick.color": MUTED,
        "xtick.labelcolor": INK_2,
        "ytick.labelcolor": INK_2,
        "grid.color": GRID,
        "grid.linewidth": 0.8,
        "axes.grid": False,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "legend.frameon": False,
        "figure.dpi": 150,
    })
    return plt


def _cmaps():
    from matplotlib.colors import LinearSegmentedColormap
    seq = LinearSegmentedColormap.from_list("rt_seq", SEQ_BLUE)
    div = LinearSegmentedColormap.from_list(
        "rt_div", [DIV_LOW, DIV_MID, DIV_HIGH])
    return seq, div


# --------------------------------------------------------------------------- #
# Small helpers
# --------------------------------------------------------------------------- #
def _fmt(value, digits=1, dash="--"):
    """Format a number for display, rendering NaN/None as an em-dash."""
    if value is None:
        return dash
    try:
        f = float(value)
    except (TypeError, ValueError):
        return str(value)
    if not np.isfinite(f):
        return dash
    return f"{f:,.{digits}f}"


def _map_axes(ax, network):
    """Configure an axis for a lon/lat map: correct aspect, no chrome."""
    lats = [n[1] for n in network.graph.nodes()]
    mid = float(np.mean(lats)) if lats else 0.0
    # 1 degree of longitude is cos(lat) as long as 1 degree of latitude; without
    # this correction every map is horizontally stretched.
    ax.set_aspect(1.0 / max(np.cos(np.radians(mid)), 1e-6))
    ax.set_xticks([])
    ax.set_yticks([])
    for spine in ax.spines.values():
        spine.set_visible(False)


def _draw_network(ax, network, values=None, cmap=None, vmin=None, vmax=None,
                  base_color=GRID, lw=1.6, base_lw=0.9):
    """Draw edges, optionally coloured by ``values`` (edge_id -> float).

    Edges with no value are drawn in the recessive base colour first, so the
    unmeasured extent of the network stays visible instead of silently
    disappearing -- a map that hides its own coverage gaps overstates the study.
    """
    from matplotlib.collections import LineCollection
    plain, coloured, cvals = [], [], []
    for eid in network.edge_ids:
        coords = network.edge_coords_lonlat(int(eid))
        v = None if values is None else values.get(int(eid))
        if v is None or not np.isfinite(v):
            plain.append(coords)
        else:
            coloured.append(coords)
            cvals.append(v)
    if plain:
        ax.add_collection(LineCollection(plain, colors=base_color,
                                         linewidths=base_lw, zorder=1))
    sm = None
    if coloured:
        lc = LineCollection(coloured, cmap=cmap, linewidths=lw, zorder=2)
        lc.set_array(np.asarray(cvals))
        if vmin is not None:
            lc.set_clim(vmin, vmax)
        sm = ax.add_collection(lc)
    ax.autoscale_view()
    all_coords = plain + coloured
    if all_coords:
        xs = [p[0] for c in all_coords for p in c]
        ys = [p[1] for c in all_coords for p in c]
        pad_x = (max(xs) - min(xs)) * 0.04 or 0.001
        pad_y = (max(ys) - min(ys)) * 0.04 or 0.001
        ax.set_xlim(min(xs) - pad_x, max(xs) + pad_x)
        ax.set_ylim(min(ys) - pad_y, max(ys) + pad_y)
    return sm


def _colorbar(fig, sm, ax, label):
    cb = fig.colorbar(sm, ax=ax, fraction=0.035, pad=0.02)
    cb.set_label(label, color=INK_2, fontsize=9)
    cb.outline.set_visible(False)
    cb.ax.tick_params(labelsize=8, color=MUTED, labelcolor=INK_2)
    return cb


def _runs(sorted_ints):
    """Yield (start, end) of each run of consecutive integers."""
    if not sorted_ints:
        return
    start = prev = sorted_ints[0]
    for h in sorted_ints[1:]:
        if h == prev + 1:
            prev = h
            continue
        yield start, prev
        start = prev = h
    yield start, prev


def build_tables(args, R):
    """The numbers behind the figures, as plain data.

    Kept format-neutral so a renderer can lay them out natively -- a real
    PowerPoint table rather than a picture of one. Every figure that encodes a
    value in colour needs its numbers readable some other way; that is the
    whole point of carrying these alongside the maps.
    """
    net, unit, out = R["net"], args.unit, {}
    if R.get("movers") is not None:
        mv = R["movers"]
        total = int(mv["n_intervals"].sum()) or 1
        rows = []
        for m in MODE_ORDER:
            sub = mv[mv["mode"] == m]
            if not len(sub):
                continue
            n_iv = int(sub["n_intervals"].sum())
            rows.append((m.capitalize(), f"{len(sub):,}", f"{n_iv:,}",
                         f"{100.0 * n_iv / total:.0f}%"))
        out["modes"] = ("Feed composition by mover mode",
                        ["Mode", "Movers", "Intervals", "Share of intervals"], rows)
    h = R["hourly"]
    if len(h):
        col = "median_speed" if "median_speed" in h.columns else "mean_speed"
        out["hourly"] = (
            "Network-wide speed by hour of day",
            ["Hour", f"Speed ({unit})", "n"],
            [(f"{int(r.block_start_hour):02d}:00", _fmt(getattr(r, col)),
              f"{int(r.n):,}")
             for r in h.sort_values("block_start_hour").itertuples()])
    if R["congestion"]:
        e = R["congestion"]["edges"]
        rated = _by_road(net, e[e["ratio"].notna()]).sort_values("ratio")
        rows = []
        for r in rated.head(12).itertuples():
            d = net.edge_data(int(r.edge_id))
            rows.append((str(d.get("name") or d.get("highway") or r.edge_id),
                         _fmt(r.observed_speed), _fmt(r.speed_limit),
                         f"{r.ratio:.2f}", f"{int(r.n):,}"))
        out["worst"] = ("Most congested roads (slower direction)",
                        ["Road", f"Observed ({unit})", f"Limit ({unit})",
                         "Ratio", "n"], rows)
    if R["bc"]:
        # Betweenness is per directed edge, so a two-way road appears twice
        # with near-identical scores; collapse to the physical road and keep
        # its busier direction, as the congestion ranking does.
        best = {}
        for eid, v in R["bc"].items():
            rid = min(net.road_edge_ids(int(eid)))
            if rid not in best or v > best[rid][1]:
                best[rid] = (int(eid), float(v))
        rows = []
        for eid, v in sorted(best.values(), key=lambda kv: -kv[1])[:12]:
            d = net.edge_data(int(eid))
            rows.append((str(d.get("name") or d.get("highway") or eid),
                         f"{v:.4f}", _fmt(d.get("length_m"), 0)))
        out["chokepoints"] = ("Highest-betweenness segments",
                              ["Segment", "Betweenness", "Length (m)"], rows)
    return out


def build_tiles(args, R):
    """Headline numbers per section, as (value, label, note) triples."""
    net, unit = R["net"], args.unit
    seg, stats, conn = R["seg"], R["stats"], R["conn"]
    pdf, out = R["points"], {}
    t0, t1 = pdf["time"].min(), pdf["time"].max()
    n_mov = int(pdf["traj_id"].nunique()) if "traj_id" in pdf.columns else 0
    out["cover"] = [
        (f"{len(pdf):,}", "GPS fixes", ""), (f"{n_mov:,}", "Movers", ""),
        (f"{len(R['intervals']):,}", "Speed intervals", "independent"),
        (f"{max((t1 - t0).days + 1, 1):,}", "Days observed", "")]
    if R.get("movers") is not None:
        mv = R["movers"]
        kept = int((mv["mode"] == "vehicle").sum())
        thr = R["mode_threshold"]
        out["modes"] = [
            (f"{kept:,}", "Movers kept", "classified as vehicles"),
            (f"{len(mv) - kept:,}", "Movers excluded", "pedestrian or unclassified"),
            (f"{thr:.1f}" if thr is not None else "--", f"Screen ({unit})",
             "per-mover 85th percentile"),
            (f"{100.0 * kept / max(len(mv), 1):.0f}%", "Of all movers", "")]
    obs = seg["coverage"]["overall"]
    out["coverage"] = [
        (f"{obs:,}", "Edges observed", ""),
        (f"{100.0 * obs / max(net.number_of_edges(), 1):.0f}%", "Of network", ""),
        (f"{seg['coverage']['peak']:,}", "In peak window", ""),
        (f"{seg['coverage']['offpeak']:,}", "In off-peak window", "")]
    ov = R["daytype"]["overall"]
    tiles = [(_fmt(ov.get(f"{lb}_speed")), f"{lb} ({unit})",
              f"n = {R['daytype']['groups'][lb]['n']:,}")
             for lb in R["daytype"]["groups"]]
    delta = ov.get("delta_pct")
    if delta is not None and np.isfinite(delta):
        tiles.append((f"{delta:+.1f}%", "Weekend vs weekday",
                      "positive = weekends faster"))
    out["daytype"] = tiles
    if R["congestion"]:
        c = R["congestion"]["summary"]
        e = R["congestion"]["edges"]
        rated = e[e["ratio"].notna()]
        out["congestion_map"] = [
            (f"{c['median_ratio']:.2f}", "Median ratio", "network-wide"),
            (f"{c['n_edges_rated']:,}", "Segments rated", "observed + posted"),
            (f"{int((rated['ratio'] < 0.5).sum()):,}", "Below 50% of limit", ""),
            (f"{int((rated['ratio'] > 1.0).sum()):,}", "Above the limit", "")]
    scc = conn.get("largest_scc_frac")
    structure = [
        (f"{stats['n_intersections']:,}", "Intersections", ""),
        (f"{stats.get('n_dead_ends', 0):,}", "Dead ends", ""),
        (_fmt(stats.get("streets_per_node_avg"), 2), "Streets per node", ""),
        (_fmt(stats.get("circuity_avg"), 3), "Circuity", "1.00 = straight"),
        ("Yes" if conn.get("is_strongly_connected") else "No",
         "Fully drivable", "every segment reaches every other"),
        (f"{100.0 * scc:.1f}%" if scc is not None else "--",
         "Largest connected part", "")]
    if stats.get("edge_density_km2") is not None:
        structure.insert(4, (_fmt(stats["edge_density_km2"], 1),
                             "Edges / km2", ""))
    out["structure"] = structure
    return out


# --------------------------------------------------------------------------- #
# Pipeline
# --------------------------------------------------------------------------- #
def run_pipeline(args, notes: list) -> dict:
    """Load, match, derive, clean and aggregate. Returns everything the deck needs."""
    print(f"[1/8] network  <- {args.network}", file=sys.stderr)
    net = rt.Network.from_geojson(args.network)
    print(f"      {net.number_of_nodes():,} nodes / {net.number_of_edges():,} edges",
          file=sys.stderr)

    print(f"[2/8] points   <- {args.points}", file=sys.stderr)
    pts = rt.load_points(args.points, tz=args.tz, id_col=args.id_col,
                         time_col=args.time_col, lon_col=args.lon_col,
                         lat_col=args.lat_col)
    pdf = pts.df
    if not pts.has_traj:
        raise SystemExit(
            "HMM matching needs a trajectory/mover id column to order fixes "
            "within a track. Pass --id-col explicitly."
        )
    if "speed_mps" in pdf.columns:
        notes.append(
            "The input carried a speed column; it was ignored. Speeds here are "
            "reconstructed from on-road displacement after matching, which is "
            "robust to the receiver's own instantaneous-speed estimate.")

    acc = args.accuracy_col
    if acc and acc in pdf.columns:
        a = pd.to_numeric(pdf[acc], errors="coerce")
        med = float(a.median()) if a.notna().any() else float("nan")
        print(f"      accuracy '{acc}': median {med:.1f} m", file=sys.stderr)
        # A horizontal accuracy in metres is typically ~3-50 m. A median below
        # ~1.5 is the signature of an HDOP/dimensionless column, which would be
        # silently interpreted as metres and make the error model far too
        # confident -- inflating the share of intervals that pass the quality
        # screen. Surface it rather than let it bias the study invisibly.
        if np.isfinite(med) and med < 1.5:
            notes.append(
                f"Accuracy column '{acc}' has a median of {med:.2f}, which is "
                "low for a horizontal accuracy in metres and typical of an HDOP "
                "(dimensionless) field. It is being used as a per-point sigma in "
                "metres; if it is HDOP, multiply it by the receiver's UERE "
                "(commonly 3-5 m) before running this report.")
    elif acc:
        notes.append(f"Accuracy column '{acc}' not found; every point fell back "
                     f"to the default sigma of {args.default_sigma} m.")
        acc = None

    print("[3/8] matching (HMM/Viterbi)", file=sys.stderr)
    matched = rt.HMMMatcher(net, max_dist=args.max_dist).match(pts)
    n_unmatched = int((matched["edge_id"] == -1).sum())
    if n_unmatched:
        pct = 100.0 * n_unmatched / max(len(matched), 1)
        notes.append(f"{n_unmatched:,} of {len(matched):,} fixes ({pct:.1f}%) "
                     f"snapped to no road within {args.max_dist} m and are "
                     "excluded from every statistic below.")

    print("[4/8] deriving speeds from on-road displacement", file=sys.stderr)
    derived = rt.derive_speeds(
        net, matched, pts,
        pos_accuracy_col=acc,
        default_pos_sigma_m=args.default_sigma,
        min_baseline_m=args.min_baseline,
    )
    intervals, obs = derived["intervals"], derived["edge_observations"]
    if not len(intervals):
        raise SystemExit(
            "No speed intervals could be derived. Common causes: fixes too far "
            "apart in time (see --max-dt), only one fix per mover, or a "
            "trajectory id that is not actually grouping the tracks."
        )
    n_bad = int((~intervals["quality"]).sum())
    bad_frac = n_bad / max(len(intervals), 1)
    if n_bad:
        notes.append(
            f"{n_bad:,} of {len(intervals):,} derived intervals "
            f"({100.0 * bad_frac:.1f}%) failed the quality screen "
            "-- displacement not clearly above the GPS noise floor, an implausible "
            "implied speed, or a poor snap. "
            + ("They are excluded (--require-quality)." if args.require_quality
               else "They are INCLUDED; re-run with --require-quality to exclude them."))
    # The quality screen is not neutral with respect to speed. It rejects
    # intervals whose displacement is small relative to GPS noise, and slow
    # traffic covers the least ground per fix -- so excluding failures
    # preferentially deletes the slowest measurements and biases every speed
    # UPWARD, understating exactly the congestion a study like this exists to
    # measure. Measured on a synthetic set with known ground truth, peak-hour
    # congestion was understated by 18% with --require-quality alone, and by
    # only 3% once --min-baseline merged hops to lift displacement clear of the
    # noise floor (which also lifted the pass rate from 75% to 95%).
    if args.require_quality and args.min_baseline is None and bad_frac > 0.10:
        notes.append(
            "WARNING -- possible upward speed bias. --require-quality is on, "
            f"{100.0 * bad_frac:.0f}% of intervals were rejected, and no "
            "--min-baseline was set. Rejected intervals are disproportionately "
            "the slow ones, so congestion here is likely UNDERSTATED. Re-run "
            "with --min-baseline 150 (or larger) to merge hops until the "
            "displacement clears the noise floor, then compare the two runs.")

    movers, mode_threshold = None, None
    if args.mode_threshold is not None:
        print("[5/8] mode screening", file=sys.stderr)
        if args.mode_threshold == "auto":
            feat = rt.mover_features(intervals, unit=args.unit)
            thr = rt.suggest_mode_threshold(feat[f"speed_p85_{args.unit}"],
                                            unit=args.unit)
            if thr is None:
                raise SystemExit(
                    "--mode-threshold auto found no walking-speed population to "
                    "split off. Either the feed is all vehicles (nothing to "
                    "exclude), or the populations overlap too much to separate on "
                    "speed. Inspect it with scripts/mover_screen.py, then pass an "
                    "explicit number.")
        else:
            thr = float(args.mode_threshold)
        mode_threshold = float(thr)

        movers = rt.classify_movers(intervals, threshold=thr, unit=args.unit)
        keep = ("vehicle", "unknown") if args.keep_unknown else ("vehicle",)
        n_before, mov_before = len(obs), int(intervals["traj_id"].nunique())
        obs = rt.filter_by_mode(obs, movers, keep=keep)
        intervals = rt.filter_by_mode(intervals, movers, keep=keep)
        kept = int(movers["mode"].isin(keep).sum())
        notes.append(
            f"Mode screening excluded {mov_before - kept:,} of {mov_before:,} "
            f"movers and {n_before - len(obs):,} of {n_before:,} observations as "
            "pedestrians or other non-vehicle movers. Mode is judged per mover on "
            "its 85th-percentile speed, so a vehicle kept by the screen retains "
            "all of its slow observations.")
        notes.append(
            "WARNING -- mode screening biases speeds UPWARD when it is wrong. A "
            "vehicle whose entire track is gridlocked never shows a fast stretch "
            "and is excluded with the pedestrians. Re-run without "
            "--mode-threshold and compare peak speeds; the gap is the uncertainty.")

    print("[6/8] cleaning", file=sys.stderr)
    # filter_trajectory_speed is deliberately NOT used here: it is the cleaner
    # for matched *point* frames (it needs lon/lat/traj_id to detect dwells),
    # whereas derive_speeds emits interval observations that carry no position.
    # The dwell case it exists to catch is already covered on this path -- a
    # parked vehicle produces displacement below the GPS noise floor, which the
    # quality screen flags -- so --require-quality is the equivalent control.
    before = len(obs)
    clean = rt.filter_by_speed(obs, max_speed=args.max_speed, unit=args.unit,
                               mad_outliers=True, per_edge=True)
    dropped = before - len(clean)
    if dropped:
        notes.append(
            f"Cleaning removed {dropped:,} of {before:,} edge observations "
            f"({100.0 * dropped / before:.1f}%) as above {args.max_speed} "
            f"{args.unit} or as robust (MAD) per-edge outliers.")
    if not len(clean):
        raise SystemExit("Cleaning removed every observation; raise --max-speed.")

    rq = args.require_quality
    print("[7/8] assigning speeds + aggregating", file=sys.stderr)
    seg = rt.assign_segment_speeds(net, clean, statistic=args.statistic,
                                   n_peak=args.n_peak, n_offpeak=args.n_offpeak,
                                   require_quality=rq)
    hourly = rt.aggregate_speeds(clean, block_hours=1, statistic="both",
                                 output_unit=args.unit, require_quality=rq,
                                 min_samples=args.min_samples)
    peaks = (rt.peak_analysis(hourly, statistic=args.statistic,
                              n_peak=args.n_peak, n_offpeak=args.n_offpeak)
             if len(hourly) else None)

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        daytype = rt.day_type_report(clean, statistic=args.statistic,
                                     output_unit=args.unit, require_quality=rq)
        for w in caught:
            if "no observations" in str(w.message):
                notes.append("One day-type had no observations, so the "
                             "weekday/weekend comparison is partial.")

    congestion = rt.congestion_report(net, clean, statistic=args.statistic,
                                      output_unit=args.unit, require_quality=rq)
    if congestion["summary"]["n_edges_rated"] == 0:
        congestion = None
        notes.append("No edge carried a usable OSM 'maxspeed' tag, so the "
                     "congestion-vs-posted-limit section is omitted. Re-export "
                     "the network with maxspeed to enable it.")

    print("[8/8] network structure", file=sys.stderr)
    stats = rt.network_stats(net, area_km2=args.area_km2)
    conn = rt.connectivity_report(net)
    try:
        bc = rt.edge_betweenness_centrality(net, weight="travel_time_s")
    except ValueError as exc:
        bc = None
        notes.append(f"Chokepoint analysis skipped: {exc}".split(" Fix by")[0])

    return {
        "net": net, "points": pdf, "matched": matched, "intervals": intervals,
        "clean": clean, "seg": seg, "hourly": hourly, "peaks": peaks,
        "daytype": daytype, "congestion": congestion, "stats": stats,
        "conn": conn, "bc": bc, "movers": movers, "mode_threshold": mode_threshold,
    }


# --------------------------------------------------------------------------- #
# Figures
# --------------------------------------------------------------------------- #
def fig_coverage(plt, net, seg):
    fig, ax = plt.subplots(figsize=(9, 6.2))
    observed = {}
    for eid in net.edge_ids:
        d = net.edge_data(int(eid))
        if d.get("obs_speed_mps_overall") is not None or \
           d.get("obs_speed_mps") is not None:
            observed[int(eid)] = 1.0
    _draw_network(ax, net, values=observed, cmap=None, base_color=GRID, lw=1.8)
    from matplotlib.collections import LineCollection
    seen = [net.edge_coords_lonlat(e) for e in observed]
    if seen:
        ax.add_collection(LineCollection(seen, colors=SERIES_1,
                                         linewidths=1.8, zorder=3))
    _map_axes(ax, net)
    cov = 100.0 * len(observed) / max(net.number_of_edges(), 1)
    ax.set_title(f"Survey coverage — {len(observed):,} of "
                 f"{net.number_of_edges():,} edges observed ({cov:.0f}%)",
                 fontsize=12, pad=12, loc="left")
    from matplotlib.lines import Line2D
    ax.legend(handles=[Line2D([], [], color=SERIES_1, lw=2, label="Observed"),
                       Line2D([], [], color=GRID, lw=2, label="No data")],
              loc="lower right", fontsize=9)
    return fig


def fig_speed_map(plt, net, period, unit, vlim=None):
    seq, _ = _cmaps()
    attr = ("obs_speed_mps_overall" if period == "overall"
            else f"obs_speed_mps_{period}")
    vals = {}
    for eid in net.edge_ids:
        v = net.edge_data(int(eid)).get(attr)
        if v is None:
            v = net.edge_data(int(eid)).get("obs_speed_mps")
        if v is not None and np.isfinite(v):
            vals[int(eid)] = float(rt.from_mps(v, unit))
    if not vals:
        return None, None
    lo, hi = (vlim if vlim else (min(vals.values()), max(vals.values())))
    fig, ax = plt.subplots(figsize=(9, 6.2))
    sm = _draw_network(ax, net, values=vals, cmap=seq, vmin=lo, vmax=hi, lw=2.0)
    _map_axes(ax, net)
    if sm:
        _colorbar(fig, sm, ax, f"Speed ({unit})")
    ax.set_title(f"{period.capitalize()} speed", fontsize=12, pad=12, loc="left")
    return fig, (lo, hi)


def fig_hourly(plt, hourly, unit, peaks):
    col = "median_speed" if "median_speed" in hourly.columns else "mean_speed"
    h = hourly.sort_values("block_start_hour")
    x = h["block_start_hour"].to_numpy()
    y = h[col].to_numpy()
    fig, ax = plt.subplots(figsize=(9, 4.4))
    ax.grid(axis="y", zorder=0)
    if {"ci95_low", "ci95_high"} <= set(h.columns) and h["ci95_low"].notna().any():
        ax.fill_between(x, h["ci95_low"], h["ci95_high"], color=SERIES_1,
                        alpha=0.15, linewidth=0, zorder=2,
                        label="95% CI of the mean")
    ax.plot(x, y, color=SERIES_1, lw=2, zorder=3, marker="o", ms=4,
            label=f"{col.split('_')[0].capitalize()} speed")
    # Direct-label only the extremes -- a number on every point is unreadable.
    if len(y):
        i_lo, i_hi = int(np.argmin(y)), int(np.argmax(y))
        for i, va in ((i_lo, "top"), (i_hi, "bottom")):
            ax.annotate(f"{y[i]:.0f}", (x[i], y[i]), textcoords="offset points",
                        xytext=(0, -14 if va == "top" else 10), ha="center",
                        fontsize=9, color=INK, fontweight="bold")
    if peaks:
        # Merge consecutive peak hours into one span. Drawing them per-hour
        # leaves a visible seam between neighbours that reads as two separate
        # windows when it is really one continuous rush period.
        hrs = sorted(r["block_start_hour"] for r in peaks["peak"])
        for run_start, run_end in _runs(hrs):
            ax.axvspan(run_start - 0.5, run_end + 0.5, color=CRITICAL,
                       alpha=0.08, lw=0, zorder=1)
    ax.set_xticks(range(0, 24, 2))
    ax.set_xlim(-0.6, 23.6)
    ax.set_xlabel("Hour of day (local clock)", fontsize=10)
    ax.set_ylabel(f"Speed ({unit})", fontsize=10)
    # Below the axes. The curve's shape changes with the data, so any in-axes
    # corner eventually sits on a line, and the top-right strip collides with
    # the title once the title is more than a few words.
    ax.legend(fontsize=9, loc="upper center", bbox_to_anchor=(0.5, -0.17),
              ncol=2)
    ax.set_title("Speed by hour of day", fontsize=12, pad=12, loc="left")
    return fig


def fig_daytype(plt, daytype, unit):
    comp = daytype["comparison"]
    labels = [k for k in daytype["groups"]]
    cols = [f"{lb}_speed" for lb in labels]
    if comp is None or not len(comp) or not all(c in comp.columns for c in cols):
        return None
    fig, ax = plt.subplots(figsize=(9, 4.4))
    ax.grid(axis="y", zorder=0)
    palette = [SERIES_1, SERIES_2]
    for i, (lb, c) in enumerate(zip(labels, cols)):
        sub = comp[["block_start_hour", c]].dropna()
        if not len(sub):
            continue
        ax.plot(sub["block_start_hour"], sub[c], lw=2, marker="o", ms=4,
                color=palette[i % 2], label=lb.capitalize(), zorder=3)
        last = sub.iloc[-1]
        ax.annotate(lb.capitalize(), (last["block_start_hour"], last[c]),
                    textcoords="offset points", xytext=(6, 0), fontsize=9,
                    color=INK_2, va="center")
    ax.set_xticks(range(0, 24, 2))
    ax.set_xlim(-0.6, 25.5)
    ax.set_xlabel("Hour of day (local clock)", fontsize=10)
    ax.set_ylabel(f"Speed ({unit})", fontsize=10)
    ax.legend(fontsize=9, loc="lower right")
    ax.set_title("Weekday vs weekend", fontsize=12, pad=12, loc="left")
    return fig


def fig_modes(plt, movers, threshold, unit, percentile=85.0):
    """Feed composition, and the distribution the screen actually cut on."""
    if movers is None or not len(movers):
        return None
    pct_col = f"speed_p{percentile:g}_{unit}"
    if pct_col not in movers.columns:
        return None
    fig, (ax0, ax1) = plt.subplots(1, 2, figsize=(9, 4.0),
                                   gridspec_kw={"width_ratios": [1, 2]})

    counts = [int((movers["mode"] == m).sum()) for m in MODE_ORDER]
    present = [(m, c) for m, c in zip(MODE_ORDER, counts) if c]
    ax0.barh([m for m, _ in present], [c for _, c in present],
             color=[MODE_COLORS[m] for m, _ in present])
    ax0.set_xlabel("Movers")
    ax0.invert_yaxis()
    for i, (_, c) in enumerate(present):
        ax0.text(c, i, f" {c:,}", va="center", fontsize=8)
    ax0.spines[["top", "right"]].set_visible(False)

    speeds = movers[pct_col].to_numpy(dtype=float)
    speeds = speeds[np.isfinite(speeds)]
    hi = float(np.percentile(speeds, 98)) if len(speeds) else 1.0
    # Neutral achromatic grey, not a series/categorical color: this histogram
    # pools every mode, and a hue here (especially one from CAT_1's blue
    # family) would read as "vehicles only" right next to the left panel's
    # blue = vehicle bar.
    ax1.hist(speeds, bins=np.linspace(0, max(hi, 1.0), 40),
             color=INK_2, edgecolor="none")
    if threshold is not None:
        ax1.axvline(threshold, color=INK, lw=1.4, ls="--")
        ax1.annotate(f"screen at {threshold:.1f} {unit}",
                     xy=(threshold, ax1.get_ylim()[1]),
                     xytext=(4, -10), textcoords="offset points",
                     fontsize=8, color=INK, ha="left", va="top")
    ax1.set_xlabel(f"Per-mover {percentile:g}th-percentile speed ({unit})")
    ax1.set_ylabel("Movers")
    ax1.spines[["top", "right"]].set_visible(False)
    fig.tight_layout()
    return fig


def fig_congestion_map(plt, net, congestion):
    _, div = _cmaps()
    e = congestion["edges"]
    rated = e[e["ratio"].notna()]
    if not len(rated):
        return None
    vals = dict(zip(rated["edge_id"].astype(int), rated["ratio"].astype(float)))
    fig, ax = plt.subplots(figsize=(9, 6.2))
    # Centre the diverging scale on 1.0 (running at the posted limit) with
    # symmetric arms, so equal visual distance means equal deviation. The span
    # comes from the 5th/95th percentile rather than the extremes: one freak
    # segment would otherwise compress every other road into the midpoint.
    # Values beyond the span saturate at the pole colour, which is the intended
    # reading ("as congested as the scale goes").
    v = np.asarray(list(vals.values()))
    span = max(abs(np.percentile(v, 5) - 1.0),
               abs(np.percentile(v, 95) - 1.0), 0.15)
    sm = _draw_network(ax, net, values=vals, cmap=div,
                       vmin=1 - span, vmax=1 + span, lw=2.0)
    _map_axes(ax, net)
    if sm:
        cb = _colorbar(fig, sm, ax, "Observed ÷ posted limit")
        cb.ax.axhline(1.0, color=INK_2, lw=1)
    ax.set_title("Congestion relative to the posted speed limit",
                 fontsize=12, pad=12, loc="left")
    return fig


def _by_road(net, edges):
    """Collapse directed edges to physical roads, keeping the worst direction.

    A two-way street is two directed edges, so ranking on edge_id lists every
    street twice -- which reads as a bug to anyone looking at the deck. Group
    on the canonical road id (the same ``min(road_edge_ids)`` idiom the rest of
    the package uses) and keep the slower direction, since this is a ranking of
    worst-case congestion.
    """
    out = edges.copy()
    out["road_id"] = [min(net.road_edge_ids(int(e))) for e in out["edge_id"]]
    idx = out.groupby("road_id")["ratio"].idxmin()
    return out.loc[idx.dropna()]


def fig_worst_corridors(plt, net, congestion, top=12):
    e = congestion["edges"]
    rated = _by_road(net, e[e["ratio"].notna()]).sort_values("ratio").head(top)
    if not len(rated):
        return None
    names = []
    for eid in rated["edge_id"].astype(int):
        d = net.edge_data(int(eid))
        names.append(str(d.get("name") or d.get("highway") or f"edge {eid}")[:34])
    fig, ax = plt.subplots(figsize=(9, max(3.0, 0.38 * len(rated) + 1.4)))
    ax.grid(axis="x", zorder=0)
    y = np.arange(len(rated))
    # One series -> one colour. A value-ramp here would double-encode length.
    ax.barh(y, rated["ratio"], color=SERIES_1, height=0.62, zorder=3)
    ax.axvline(1.0, color=BASELINE, lw=1.2, zorder=4)
    ax.text(1.0, -0.9, "posted limit", fontsize=8, color=MUTED, ha="center")
    for i, v in enumerate(rated["ratio"]):
        ax.text(v + 0.02, i, f"{v:.2f}", va="center", fontsize=9, color=INK)
    ax.set_yticks(y, names, fontsize=9)
    ax.invert_yaxis()
    ax.set_xlim(0, max(1.15, float(rated["ratio"].max()) * 1.18))
    ax.set_xlabel("Observed speed ÷ posted limit (slower direction)", fontsize=10)
    ax.set_title(f"Most congested roads (worst {len(rated)})",
                 fontsize=12, pad=12, loc="left")
    return fig


def fig_chokepoints(plt, net, bc):
    seq, _ = _cmaps()
    vals = {int(k): float(v) for k, v in bc.items() if np.isfinite(v)}
    if not vals or max(vals.values()) <= 0:
        return None
    fig, ax = plt.subplots(figsize=(9, 6.2))
    sm = _draw_network(ax, net, values=vals, cmap=seq, vmin=0,
                       vmax=max(vals.values()), lw=2.0)
    _map_axes(ax, net)
    if sm:
        _colorbar(fig, sm, ax, "Edge betweenness")
    ax.set_title("Chokepoints — travel-time-weighted edge betweenness",
                 fontsize=12, pad=12, loc="left")
    return fig


# --------------------------------------------------------------------------- #
# Section headings, in one place so the HTML deck, the PDF and the PNG captions
# cannot drift apart. ``sub`` is a template filled from the context built in
# :func:`_section_ctx`.
# --------------------------------------------------------------------------- #
SECTIONS = [
    ("coverage", "01 — Survey coverage", "Where the data actually is",
     "Roads in grey were never traversed by the survey and carry no measured "
     "speed. Every statistic in this report describes the blue extent only."),
    ("modes", "02 — Feed composition", "What the feed is made of",
     "Movers are classified by their own 85th-percentile speed, not by "
     "individual observations: a pedestrian is slow for a whole track, while "
     "a congested vehicle is slow here and free-flowing elsewhere. A vehicle "
     "kept by the screen keeps every one of its slow observations."),
    ("speed_overall", "03 — Measured speed", "Overall speed by segment",
     "The {statistic} speed across the whole survey period. Darker is faster."),
    ("speed_peak", "04 — Peak vs off-peak",
     "The same roads, at their worst and their best",
     "Both maps share one colour scale, so they are directly comparable. Peak "
     "hours ({peak_h}) were chosen from the data as the slowest contiguous "
     "window; off-peak ({off_h}) the fastest."),
    ("speed_offpeak", "04 — Peak vs off-peak (continued)",
     "Off-peak speed by segment",
     "The fastest contiguous window, on the same scale as the peak map."),
    ("hourly", "05 — Time of day", "When the network slows down",
     "Network-wide {statistic} speed by hour. Each interval counts once "
     "regardless of how many segments it crossed, so the sample sizes are "
     "independent measurements."),
    ("daytype", "06 — Day type",
     "Weekday and weekend traffic are different populations",
     "Pooling them into one hour-of-day average hides both. Split here so a "
     "Tuesday 09:00 is never averaged with a Saturday 09:00."),
    ("congestion_map", "07 — Congestion vs posted limit",
     "How the network performs against its own speed limits",
     "A ratio of 1.00 means traffic moved at the posted limit; 0.45 means it "
     "crawled at 45% of it. This is the comparison raw speeds cannot make."),
    ("worst", "07 — Congestion vs posted limit (continued)",
     "Worst-performing roads",
     "One row per physical road, showing its slower direction."),
    ("chokepoints", "08 — Network structure", "Chokepoints",
     "Edge betweenness weighted by measured travel time: the share of fastest "
     "routes that must cross each segment. This finds roads that are "
     "structurally load-bearing, not merely central on a map."),
]
SECTION_META = {k: (e, h, sb) for k, e, h, sb in SECTIONS}


def _section_ctx(args, R):
    """Values the section ``sub`` templates interpolate."""
    seg = R["seg"]
    return {
        "statistic": args.statistic,
        "peak_h": ", ".join(f"{h:02d}:00" for h in seg["peak_hours"][:6]),
        "off_h": ", ".join(f"{h:02d}:00" for h in seg["offpeak_hours"][:6]),
    }


# --------------------------------------------------------------------------- #
# Figures, then one renderer per output format
# --------------------------------------------------------------------------- #
def build_figures(plt, R, args) -> dict:
    """Every figure the report can show, as live matplotlib figures.

    Returned open rather than serialised so each renderer can decide what to do
    with them -- base64 for the HTML deck, a page for the PDF, a file for the
    PNG bundle -- without the analysis being run three times. Sections whose
    data is missing are simply absent from the dict, and every renderer keys off
    that rather than testing conditions again.
    """
    figs, net = {}, R["net"]
    figs["coverage"] = fig_coverage(plt, net, R["seg"])
    if R.get("movers") is not None:
        f = fig_modes(plt, R["movers"], R["mode_threshold"], args.unit)
        if f:
            figs["modes"] = f
    f, lim = fig_speed_map(plt, net, "overall", args.unit)
    if f:
        figs["speed_overall"] = f
    # Peak and off-peak share the overall scale so the maps compare directly.
    for period in ("peak", "offpeak"):
        f, _ = fig_speed_map(plt, net, period, args.unit, vlim=lim)
        if f:
            figs[f"speed_{period}"] = f
    if len(R["hourly"]):
        figs["hourly"] = fig_hourly(plt, R["hourly"], args.unit, R["peaks"])
    f = fig_daytype(plt, R["daytype"], args.unit)
    if f:
        figs["daytype"] = f
    if R["congestion"]:
        f = fig_congestion_map(plt, net, R["congestion"])
        if f:
            figs["congestion_map"] = f
        f = fig_worst_corridors(plt, net, R["congestion"])
        if f:
            figs["worst"] = f
    if R["bc"]:
        f = fig_chokepoints(plt, net, R["bc"])
        if f:
            figs["chokepoints"] = f
    return figs


def _pdf_text_page(plt, pdf, *, eyebrow, heading, paragraphs, tiles=None,
                   size=(11.0, 8.5)):
    """A text-only PDF page laid out to match the deck's typographic rhythm."""
    fig = plt.figure(figsize=size, facecolor=SURFACE)
    y = 0.90
    if eyebrow:
        fig.text(0.07, y, eyebrow.upper(), fontsize=9, color=MUTED,
                 fontweight="bold")
        y -= 0.045
    fig.text(0.07, y, heading, fontsize=21, color=INK, fontweight="bold")
    y -= 0.06
    if tiles:
        for i, (val, key) in enumerate(tiles):
            x = 0.07 + (i % 4) * 0.225
            row = i // 4
            fig.text(x, y - row * 0.10, str(val), fontsize=20, color=INK,
                     fontweight="bold")
            fig.text(x, y - row * 0.10 - 0.032, key.upper(), fontsize=8,
                     color=MUTED, fontweight="bold")
        y -= 0.10 * (1 + (len(tiles) - 1) // 4) + 0.02
    for para in paragraphs:
        for line in textwrap.wrap(para, width=104) or [""]:
            fig.text(0.07, y, line, fontsize=10.5, color=INK_2)
            y -= 0.028
        y -= 0.018
    pdf.savefig(fig, facecolor=SURFACE)
    plt.close(fig)


def render_pdf(args, R, live, notes, plt):
    """Multi-page PDF -- no extra dependencies, fixed and print-ready."""
    from matplotlib.backends.backend_pdf import PdfPages
    ctx = _section_ctx(args, R)
    pdf_ctx = R["points"]
    t0, t1 = pdf_ctx["time"].min(), pdf_ctx["time"].max()
    n_mov = (int(pdf_ctx["traj_id"].nunique())
             if "traj_id" in pdf_ctx.columns else 0)
    with PdfPages(args.out) as pdf:
        _pdf_text_page(
            plt, pdf, eyebrow="Trafficability study", heading=args.title,
            tiles=[(f"{len(pdf_ctx):,}", "GPS fixes"), (f"{n_mov:,}", "Movers"),
                   (f"{len(R['intervals']):,}", "Speed intervals"),
                   (f"{max((t1 - t0).days + 1, 1):,}", "Days observed")],
            paragraphs=[
                args.customer,
                f"Period {t0:%Y-%m-%d %H:%M} to {t1:%Y-%m-%d %H:%M}"
                f"  |  network {R['net'].number_of_edges():,} edges"
                f"  |  generated {datetime.now(timezone.utc):%Y-%m-%d} UTC",
                "Speeds in this report were not read from the GPS receiver. "
                "Each fix was map-matched to the road network with a hidden "
                "Markov model, and speed reconstructed from on-road "
                "displacement between consecutive fixes of the same mover.",
            ])
        for key, eyebrow, _heading, sub_t in SECTIONS:
            if key not in live:
                continue
            fig = live[key]
            # Eyebrow and caption go ABOVE the plot, anchored va="bottom" so
            # they grow upward into empty space. Below the plot is already
            # occupied on some figures (the hourly chart puts its legend
            # there), and a caption placed underneath lands on top of it.
            fig.text(0.02, 1.02, textwrap.fill(sub_t.format(**ctx), 116),
                     va="bottom", fontsize=9, color=INK_2)
            fig.text(0.02, 1.13, eyebrow.upper(), fontsize=8, color=MUTED,
                     fontweight="bold", va="bottom")
            pdf.savefig(fig, bbox_inches="tight", facecolor=SURFACE)
        _pdf_text_page(
            plt, pdf, eyebrow="10 — Method",
            heading="How these numbers were produced",
            paragraphs=[
                "Matching. Each fix was assigned to a road with a hidden Markov "
                "model decoded by Viterbi (Newson & Krumm, 2009), which uses the "
                "whole trajectory rather than snapping each point independently.",
                "Speed. Derived from on-road displacement between consecutive "
                "fixes of one mover, divided by elapsed time, with per-point GPS "
                "accuracy propagated into an explicit uncertainty per interval.",
                f"Aggregation. {args.statistic.capitalize()} statistics; each "
                "interval counts once network-wide however many segments it "
                "crossed. Peak and off-peak windows were selected from the data.",
                "Limitations. Speeds describe the surveyed vehicles, not all "
                "traffic, and only where and when they drove. A posted limit is "
                "a legal maximum, not a free-flow speed, so the congestion ratio "
                "is a screening indicator; comparing a road against itself "
                "across time blocks is sounder than comparing different roads.",
            ])
        if notes:
            _pdf_text_page(plt, pdf, eyebrow="Data notes",
                           heading="Caveats on this dataset",
                           paragraphs=[f"- {n}" for n in notes])
    return args.out


# 16:9 slide geometry, shared by the layout helpers below.
SLIDE_W, SLIDE_H = 13.333, 7.5
PIC_TOP = 2.15


def _png_size(path):
    """Pixel dimensions of a PNG, for aspect-correct placement."""
    from PIL import Image
    with Image.open(path) as im:
        return im.size


def _pptx_rgb(hex_str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hex_str.lstrip("#").upper())


def _pptx_textbox(slide, x, y, w, h, text, *, size, color, bold=False,
                  caps=False, space_after=0):
    from pptx.util import Inches, Pt
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, para_text in enumerate(text if isinstance(text, list) else [text]):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = para_text.upper() if caps else para_text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = _pptx_rgb(color)
        run.font.name = "Calibri"
        para.space_after = Pt(space_after)
    return box


def _pptx_tiles(slide, tiles, y):
    """A row of headline numbers, wrapping every four."""
    for i, (val, label, note) in enumerate(tiles):
        x = 0.55 + (i % 4) * 3.05
        row_y = y + (i // 4) * 1.05
        _pptx_textbox(slide, x, row_y, 2.9, 0.5, str(val), size=26,
                      color=INK, bold=True)
        _pptx_textbox(slide, x, row_y + 0.48, 2.9, 0.3, label, size=9,
                      color=MUTED, bold=True, caps=True)
        if note:
            _pptx_textbox(slide, x, row_y + 0.70, 2.9, 0.3, note, size=9,
                          color=INK_2)


def _pptx_table(slide, headers, rows, x, y, w, h):
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    table = shape.table
    for c, head in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(head)
        para = cell.text_frame.paragraphs[0]
        para.runs[0].font.size = Pt(10)
        para.runs[0].font.bold = True
        para.runs[0].font.color.rgb = _pptx_rgb(MUTED)
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(10)
            run.font.color.rgb = _pptx_rgb(INK)
    return shape


def render_pptx(args, R, live, notes, plt):
    """Editable PowerPoint deck -- one slide per section, native tables.

    The only format the customer can restyle or fold into their own template,
    which is the reason to pay for the extra dependency. Tables are real
    PowerPoint tables rather than pictures of tables, so every value a figure
    encodes in colour stays readable (and selectable) without it.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise SystemExit(
            "The pptx format needs python-pptx:\n    pip install python-pptx\n"
            "(or use --format pdf / --format png, which need no extra deps)"
        ) from exc
    import os
    import tempfile

    ctx = _section_ctx(args, R)
    tables, tiles = build_tables(args, R), build_tiles(args, R)
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    def new_slide(eyebrow, heading, sub=None):
        sl = prs.slides.add_slide(blank)
        sl.background.fill.solid()
        sl.background.fill.fore_color.rgb = _pptx_rgb(SURFACE)
        _pptx_textbox(sl, 0.55, 0.35, 12.2, 0.3, eyebrow, size=10,
                      color=MUTED, bold=True, caps=True)
        _pptx_textbox(sl, 0.55, 0.68, 12.2, 0.6, heading, size=26, color=INK,
                      bold=True)
        if sub:
            _pptx_textbox(sl, 0.55, 1.32, 12.2, 0.7, sub, size=12, color=INK_2)
        return sl

    pdf_pts = R["points"]
    t0, t1 = pdf_pts["time"].min(), pdf_pts["time"].max()
    cover = new_slide("Trafficability study", args.title, args.customer)
    _pptx_tiles(cover, tiles["cover"], 2.3)
    _pptx_textbox(
        cover, 0.55, 4.6, 12.2, 1.6,
        [f"Period {t0:%Y-%m-%d %H:%M} to {t1:%Y-%m-%d %H:%M}   |   "
         f"network {R['net'].number_of_edges():,} edges   |   "
         f"generated {datetime.now(timezone.utc):%Y-%m-%d} UTC",
         "Speeds were not read from the GPS receiver. Each fix was map-matched "
         "to the road network with a hidden Markov model, and speed "
         "reconstructed from on-road displacement between consecutive fixes of "
         "the same mover."],
        size=12, color=INK_2, space_after=8)

    tmp = tempfile.mkdtemp(prefix="rt_pptx_")
    for key, eyebrow, heading, sub_t in SECTIONS:
        if key not in live:
            continue
        sl = new_slide(eyebrow, heading, sub_t.format(**ctx))
        png = os.path.join(tmp, f"{key}.png")
        live[key].savefig(png, dpi=args.dpi, bbox_inches="tight",
                          facecolor=SURFACE)
        has_side = key in tables or key in tiles
        # Fit the picture to BOTH the available width and the space left above
        # the bottom of the slide. Sizing on width alone silently runs tall
        # figures -- the maps are 9x6.2in, so they scale to over 9in high --
        # straight off the bottom of a 7.5in slide.
        max_w = 8.2 if has_side else 11.4
        max_h = SLIDE_H - PIC_TOP - 0.3
        img_w, img_h = _png_size(png)
        aspect = img_h / img_w
        w = min(max_w, max_h / aspect)
        sl.shapes.add_picture(png, Inches(0.55 + (max_w - w) / 2),
                              Inches(PIC_TOP), width=Inches(w))
        if key in tables:
            caption, headers, rows = tables[key]
            _pptx_textbox(sl, 9.05, 2.15, 3.7, 0.3, caption, size=9,
                          color=MUTED, bold=True, caps=True)
            _pptx_table(sl, headers, rows[:10], 9.05, 2.5, 3.75,
                        min(0.28 * (len(rows[:10]) + 1), 4.6))
        elif key in tiles:
            for i, (val, label, note) in enumerate(tiles[key]):
                y = 2.3 + i * 1.05
                _pptx_textbox(sl, 9.05, y, 3.7, 0.4, str(val), size=22,
                              color=INK, bold=True)
                _pptx_textbox(sl, 9.05, y + 0.42, 3.7, 0.3, label, size=9,
                              color=MUTED, bold=True, caps=True)
                if note:
                    _pptx_textbox(sl, 9.05, y + 0.62, 3.7, 0.3, note, size=8,
                                  color=INK_2)

    sl = new_slide("09 — Network structure", "Structural summary",
                   "Properties of the road network itself, independent of the "
                   "GPS survey.")
    _pptx_tiles(sl, tiles["structure"], 2.3)

    sl = new_slide("10 — Method", "How these numbers were produced")
    _pptx_textbox(sl, 0.55, 1.5, 12.2, 4.6, [
        "Matching. Each fix was assigned to a road with a hidden Markov model "
        "decoded by Viterbi (Newson & Krumm, 2009), using the whole trajectory "
        "rather than snapping each point independently.",
        "Speed. Derived from on-road displacement between consecutive fixes of "
        "one mover, divided by elapsed time, with per-point GPS accuracy "
        "propagated into an explicit uncertainty for every interval.",
        f"Aggregation. {args.statistic.capitalize()} statistics; each interval "
        "counts once network-wide however many segments it crossed. Peak and "
        "off-peak windows were selected from the data.",
        "Limitations. Speeds describe the surveyed vehicles, not all traffic, "
        "and only where and when they drove. A posted limit is a legal maximum, "
        "not a free-flow speed, so the congestion ratio is a screening "
        "indicator; comparing a road against itself across time blocks is "
        "sounder than comparing different roads to each other.",
    ], size=12, color=INK_2, space_after=12)
    if notes:
        sl = new_slide("Data notes", "Caveats on this dataset")
        _pptx_textbox(sl, 0.55, 1.5, 12.2, 4.8,
                      [f"\u2022  {n}" for n in notes], size=11, color=INK_2,
                      space_after=10)

    prs.save(args.out)
    for f in os.listdir(tmp):
        os.remove(os.path.join(tmp, f))
    os.rmdir(tmp)
    return args.out


def render_png(args, R, live, notes, plt):
    """A numbered PNG per section plus a captions file, for your own template."""
    import os
    out_dir = args.out
    if out_dir.lower().endswith((".png", ".pdf", ".pptx")):
        out_dir = os.path.splitext(out_dir)[0]
    os.makedirs(out_dir, exist_ok=True)
    ctx = _section_ctx(args, R)
    lines = [f"# {args.title}", ""]
    if args.customer:
        lines += [args.customer, ""]
    n = 0
    for key, eyebrow, heading, sub_t in SECTIONS:
        if key not in live:
            continue
        n += 1
        name = f"{n:02d}_{key}.png"
        live[key].savefig(os.path.join(out_dir, name), dpi=args.dpi,
                          bbox_inches="tight", facecolor=SURFACE)
        lines += [f"## {name}", f"**{eyebrow} — {heading}**", "",
                  textwrap.fill(sub_t.format(**ctx), 92), ""]
    if notes:
        lines += ["## Data notes", ""] + [f"- {x}" for x in notes]
    with open(os.path.join(out_dir, "captions.md"), "w", encoding="utf-8") as fh:
        fh.write("\n".join(lines) + "\n")
    return f"{out_dir}/ ({n} PNGs + captions.md)"


# --------------------------------------------------------------------------- #
def parse_args(argv=None):
    p = argparse.ArgumentParser(
        description="Build a customer-ready congestion report from GPS + a road network.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""\
            example:
              python scripts/customer_report.py --network roads.geojson \\
                --points probes.csv --accuracy-col accuracy \\
                --tz America/New_York --title "Route 1 Study" \\
                --customer "Example DOT" --out report.pdf
            """))
    p.add_argument("--network", required=True, help="road network GeoJSON")
    p.add_argument("--points", required=True, help="GPS points CSV/TSV/GeoJSON")
    p.add_argument("--out", default="report.pdf",
                   help="output file (or directory, for --format png)")
    p.add_argument("--format", default=None,
                   choices=["pdf", "pptx", "png"],
                   help="deliverable format; inferred from --out's extension "
                        "when omitted (default pdf)")
    p.add_argument("--dpi", type=int, default=200,
                   help="raster resolution for --format png (default 200)")
    p.add_argument("--title", default="Road Network Trafficability Study")
    p.add_argument("--customer", default="")
    p.add_argument("--tz", default=None,
                   help="IANA timezone of the study area, e.g. America/New_York")
    p.add_argument("--accuracy-col", default="accuracy",
                   help="per-point horizontal accuracy in METRES (default: accuracy)")
    p.add_argument("--id-col", default=None, help="mover/trajectory id column")
    p.add_argument("--time-col", default=None)
    p.add_argument("--lon-col", default=None)
    p.add_argument("--lat-col", default=None)
    p.add_argument("--unit", default="mph", choices=["mph", "kph", "mps"])
    p.add_argument("--statistic", default="median", choices=["median", "mean"])
    p.add_argument("--max-dist", type=float, default=50.0,
                   help="max snap distance, metres (default 50)")
    p.add_argument("--default-sigma", type=float, default=15.0,
                   help="fallback GPS sigma when accuracy is missing (default 15 m)")
    p.add_argument("--min-baseline", type=float, default=None,
                   help="merge hops until this on-road distance; raises SNR on "
                        "dense/noisy fixes")
    p.add_argument("--max-speed", type=float, default=80.0,
                   help="physical upper bound for cleaning, in --unit")
    p.add_argument("--min-samples", type=int, default=1,
                   help="suppress time bins with fewer observations than this "
                        "(default 1; raise to stop thin hours reading as spikes)")
    p.add_argument("--n-peak", type=int, default=3)
    p.add_argument("--n-offpeak", type=int, default=3)
    p.add_argument("--area-km2", type=float, default=None,
                   help="study-area size, enables per-km2 densities")
    p.add_argument("--require-quality", action="store_true",
                   help="drop intervals that failed the quality screen "
                        "(recommended for customer deliverables)")
    p.add_argument("--mode-threshold", default=None,
                   help="exclude movers whose p85 speed is below this (in --unit), "
                        "or 'auto' to pick the density valley. Off by default: "
                        "screening changes the study, so it is never implicit.")
    p.add_argument("--keep-unknown", action="store_true",
                   help="keep movers with too little data to classify "
                        "(default: excluded along with pedestrians)")
    return p.parse_args(argv)


def main(argv=None):
    args = parse_args(argv)
    if args.format is None:
        low = args.out.lower()
        args.format = ("pptx" if low.endswith(".pptx")
                       else "png" if low.endswith(".png") else "pdf")
    notes: list = []
    R = run_pipeline(args, notes)
    plt = _mpl()

    print("      rendering figures", file=sys.stderr)
    live = build_figures(plt, R, args)
    render = {"pdf": render_pdf, "pptx": render_pptx, "png": render_png}
    out = render[args.format](args, R, live, notes, plt)
    for f in live.values():                 # renderers must not close them
        plt.close(f)
    print(f"\nWrote {out}  ({len(live)} figures, format={args.format})",
          file=sys.stderr)
    if notes:
        print(f"{len(notes)} data note(s) included in the report.", file=sys.stderr)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
